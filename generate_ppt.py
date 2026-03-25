from pptx import Presentation
from pptx.util import Inches, Pt

def create_ppt():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "vaptAI: AI-Powered Vulnerability Assessment"
    subtitle.text = "Automated Pentesting & Forensic Auditing\nMajor Project Presentation"

    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Introduction"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "vaptAI is an end-to-end security framework."
    p = tf.add_paragraph()
    p.text = "Combines industry-standard scanners with Generative AI."
    p = tf.add_paragraph()
    p.text = "Automates the VAPT lifecycle: Recon, Scan, Analysis, Reporting."

    # Slide 3: Problem Statement
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Problem Statement"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Complexity of raw security logs from tools like Nmap/Nuclei."
    p = tf.add_paragraph()
    p.text = "High barrier of entry for non-experts to interpret results."
    p = tf.add_paragraph()
    p.text = "Manual reporting is slow and prone to errors."

    # Slide 4: System Architecture
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "System Architecture"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Frontend: React.js (Vite) for interactive dashboards."
    p = tf.add_paragraph()
    p.text = "Backend: Flask (Python) with Asynchronous Threading."
    p = tf.add_paragraph()
    p.text = "Database: SQLite for persistence and audit logs."

    # Slide 5: Core Modules - Technical Scanning
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Core Module: Technical Scanning"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Nmap: Port discovery and service fingerprinting."
    p = tf.add_paragraph()
    p.text = "Nuclei: Template-based vulnerability scanning (CVEs)."
    p = tf.add_paragraph()
    p.text = "Subdomain Enumeration: Custom DNS brute-forcing."

    # Slide 6: AI Intelligence Layer
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "AI Intelligence Layer"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Powered by Google Gemini 2.0 Flash."
    p = tf.add_paragraph()
    p.text = "Context-aware analysis and risk scoring."
    p = tf.add_paragraph()
    p.text = "AI Chat Assistant for real-time remediation advice."
    p = tf.add_paragraph()
    p.text = "Offline Mode fallback for basic heuristic analysis."

    # Slide 7: Forensics & Data Integrity
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Forensics & Data Integrity"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Automated JSON evidence file generation."
    p = tf.add_paragraph()
    p.text = "SHA-256 Hashing of scan results to ensure non-repudiation."
    p = tf.add_paragraph()
    p.text = "Validation checks to detect report tampering."

    # Slide 8: OSINT Integration
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "OSINT Integration"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Integrated Sherlock for username tracking."
    p = tf.add_paragraph()
    p.text = "Maps digital footprints across 300+ platforms."
    p = tf.add_paragraph()
    p.text = "Enhances social engineering risk assessments."

    # Slide 9: Results & Reporting
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Results & Reporting"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Automated PDF Report Generation (ReportLab)."
    p = tf.add_paragraph()
    p.text = "Executive Summary for management + Technical logs for IT."
    p = tf.add_paragraph()
    p.text = "Step-by-step remediation guides for all findings."

    # Slide 10: Conclusion & Future Work
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Conclusion & Future Work"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Future: Dockerization for easier deployment."
    p = tf.add_paragraph()
    p.text = "Cloud-native scanning (AWS/Azure buckets)."
    p = tf.add_paragraph()
    p.text = "Continuous security monitoring via scheduled scans."

    prs.save('vaptAI_Project_Presentation.pptx')
    print("PPT created successfully: vaptAI_Project_Presentation.pptx")

if __name__ == "__main__":
    create_ppt()
