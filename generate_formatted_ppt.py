from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_formatted_ppt():
    prs = Presentation()

    # Slide 1: Title Slide (Based on format)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "vaptAI: AI-Powered Vulnerability Assessment and Penetration Testing Framework"
    subtitle.text = "Major Project Presentation\nStudent Name: [Your Name]\nRegistration No: [Your Reg No]\nFaculty of Computer Science and Engineering"

    # Slide 2: Contents
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Contents"
    tf = slide.shapes.placeholders[1].text_frame
    items = ["Introduction", "Problem Statement", "System Architecture", "Core Modules & Forensic Integrity", "AI & OSINT Integration", "References", "Conclusion"]
    for item in items:
        p = tf.add_paragraph()
        p.text = item

    # Slide 3: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Overview: An automated platform to simplify complex security audits."
    p = tf.add_paragraph()
    p.text = "Problem Statement: Traditional tools generate overwhelming data and lack clear remediation paths."
    p = tf.add_paragraph()
    p.text = "Importance: Bridges the gap between technical scanning and executive understanding using AI."

    # Slide 4: System Architecture (The "Empty" slide 4 from sample)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "System Architecture & Workflow"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Frontend: React.js (Vite) - Real-time progress monitoring via polling."
    p = tf.add_paragraph()
    p.text = "Backend: Flask (Python) - Handles asynchronous tool orchestration."
    p = tf.add_paragraph()
    p.text = "Storage: SQLite & Forensic Evidence Folder with SHA-256 integrity."

    # Slide 5: Core Technical Modules
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Core Technical Modules"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Scanner: Nmap (Ports) + Nuclei (CVEs) + Subdomain Discovery."
    p = tf.add_paragraph()
    p.text = "AI Helper: Google Gemini 2.0 - Analysis, Scoring, and Chat."
    p = tf.add_paragraph()
    p.text = "Forensics: Automatic hashing of results to ensure non-repudiation."

    # Slide 6: REFERENCES
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "REFERENCES"
    tf = slide.shapes.placeholders[1].text_frame
    refs = [
        "[1] Nmap Project. 'Nmap Network Scanning.' nmap.org",
        "[2] ProjectDiscovery. 'Nuclei: Fast and customizable vulnerability scanner.' github.com",
        "[3] Google. 'Gemini 2.0 Flash Documentation.' ai.google.dev",
        "[4] Sherlock Project. 'Sherlock: Find usernames across social networks.' github.com"
    ]
    for ref in refs:
        p = tf.add_paragraph()
        p.text = ref

    # Slide 7: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = "Thank You"
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Any Questions?\n\n[Your Name]\n[Your Email/Contact Information]"

    prs.save('vaptAI_Formatted_Major_Project.pptx')
    print("PPT created successfully: vaptAI_Formatted_Major_Project.pptx")

if __name__ == "__main__":
    create_formatted_ppt()
