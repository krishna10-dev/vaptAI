from pptx import Presentation

def inspect_pptx(path):
    prs = Presentation(path)
    print(f"Total Slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else "No Title"
        print(f"Slide {i+1}: {title}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                print(f"  - Content: {shape.text[:50]}...")

if __name__ == "__main__":
    inspect_pptx(r"C:\Users\krishna\Downloads\Sample PPT for Major Project (1).pptx")
